#!/usr/bin/env python3
"""
Chelestra-Sea #2 raw text extraction pipeline.
Processes files from the documents table where raw_content IS NULL.
Updates raw_content column with extracted plain text.
"""
import os
import sqlite3
import logging
import sys
import time
from pathlib import Path

# Text extraction libraries
import PyPDF2
import docx
import pptx
import openpyxl
import json
import re

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('/data/intelligence/extraction.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

DB_PATH = '/data/intelligence/owner-intelligence.db'

def connect_db():
    """Return a DB connection."""
    return sqlite3.connect(DB_PATH)

def get_files_to_process(limit=1000):
    """Return list of (id, file_path, file_type) for files needing raw_content."""
    conn = connect_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT id, file_path, file_type 
        FROM documents 
        WHERE raw_content IS NULL 
        ORDER BY id 
        LIMIT ?
    """, (limit,))
    rows = cursor.fetchall()
    conn.close()
    return rows

def extract_pdf(path):
    """Extract text from PDF."""
    text = []
    try:
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"PDF extraction error {path}: {e}")
        return None

def extract_docx(path):
    """Extract text from DOCX."""
    try:
        doc = docx.Document(path)
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        logger.error(f"DOCX extraction error {path}: {e}")
        return None

def extract_pptx(path):
    """Extract text from PPTX."""
    try:
        pres = pptx.Presentation(path)
        text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"PPTX extraction error {path}: {e}")
        return None

def extract_xlsx(path):
    """Extract text from XLSX (cell values)."""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text.append(row_text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"XLSX extraction error {path}: {e}")
        return None

def extract_txt(path):
    """Read plain text file."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"TXT read error {path}: {e}")
        return None

def extract_csv(path):
    """Read CSV as text."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"CSV read error {path}: {e}")
        return None

def extract_json(path):
    """Read JSON and convert to text."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            return json.dumps(data, indent=2)
    except Exception as e:
        logger.error(f"JSON read error {path}: {e}")
        return None

def extract_html(path):
    """Read HTML and extract text (basic)."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            # Basic HTML stripping
            clean = re.sub('<[^<]+?>', ' ', content)
            clean = re.sub('\\s+', ' ', clean)
            return clean.strip()
    except Exception as e:
        logger.error(f"HTML read error {path}: {e}")
        return None

# Supported MIME types for text extraction
SUPPORTED_TYPES = {
    'application/pdf': extract_pdf,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_docx,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_pptx,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': extract_xlsx,
    'text/plain': extract_txt,
    'text/csv': extract_csv,
    'application/json': extract_json,
    'text/html': extract_html,
    'text/xml': extract_txt,
    'text/x-script.python': extract_txt,
    'application/x-bytecode.python': extract_txt,
    'text/x-php': extract_txt,
    'text/x-shellscript': extract_txt,
}

UNSUPPORTED_PLACEHOLDER = "[UNSUPPORTED_FILE_TYPE]"

def extract_text(file_path, file_type):
    """Extract plain text from file based on type."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    # Check if type is supported
    if file_type in SUPPORTED_TYPES:
        extract_func = SUPPORTED_TYPES[file_type]
        return extract_func(file_path)
    else:
        logger.warning(f"Unsupported file type: {file_type} for {file_path}")
        return UNSUPPORTED_PLACEHOLDER

def update_raw_content(file_id, text):
    """Update raw_content column for given file_id."""
    conn = connect_db()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE documents 
        SET raw_content = ? 
        WHERE id = ?
    """, (text, file_id))
    conn.commit()
    conn.close()

def process_batch(batch_size=100):
    """Process one batch of files."""
    files = get_files_to_process(limit=batch_size)
    if not files:
        logger.info("No files left to process.")
        return 0
    
    processed = 0
    unsupported = 0
    errors = 0
    
    for file_id, file_path, file_type in files:
        try:
            logger.info(f"Processing {file_path} ({file_type})")
            text = extract_text(file_path, file_type)
            if text is not None:
                # If it's unsupported placeholder, still update to mark as processed
                update_raw_content(file_id, text)
                processed += 1
                if text == UNSUPPORTED_PLACEHOLDER:
                    unsupported += 1
                    logger.debug(f"Marked as unsupported: {file_path}")
                if processed % 20 == 0:
                    logger.info(f"Progress: {processed} files processed in this batch.")
            else:
                logger.warning(f"Could not extract text from {file_path}")
                errors += 1
        except Exception as e:
            logger.error(f"Unexpected error processing {file_path}: {e}")
            errors += 1
    
    logger.info(f"Batch complete. Processed: {processed} (unsupported: {unsupported}), Errors: {errors}")
    return processed

def main():
    """Main pipeline loop."""
    logger.info("Starting Chelestra-Sea #2 extraction pipeline")
    
    # Get total count
    conn = connect_db()
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM documents WHERE raw_content IS NULL")
    total_remaining = cursor.fetchone()[0]
    conn.close()
    
    logger.info(f"Total files remaining: {total_remaining}")
    
    batch_size = 100
    total_processed = 0
    
    while total_remaining > 0:
        processed = process_batch(batch_size)
        if processed == 0:
            break
        
        total_processed += processed
        
        # Update remaining count
        conn = connect_db()
        cursor = conn.cursor()
        cursor.execute("SELECT COUNT(*) FROM documents WHERE raw_content IS NULL")
        total_remaining = cursor.fetchone()[0]
        conn.close()
        
        logger.info(f"Overall progress: {total_processed} processed, {total_remaining} remaining")
        
        # Small pause between batches
        time.sleep(2)
    
    logger.info("Extraction pipeline completed.")

if __name__ == "__main__":
    main()