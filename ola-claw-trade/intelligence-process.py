#!/usr/bin/env python3
"""
Process files listed in file-catalogue.txt, extract text content,
and update owner-intelligence.db with raw_content.
"""

import sqlite3
import os
import sys
import logging
from pathlib import Path
import PyPDF2
import docx
import pptx
import openpyxl
import csv
import json
import traceback

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

DB_PATH = "/data/intelligence/owner-intelligence.db"
CATALOGUE_PATH = "/data/intelligence/file-catalogue.txt"

def extract_text(file_path, mime_type):
    """Extract text from file based on mime type."""
    try:
        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            return None
        
        # PDF
        if mime_type == 'application/pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text()
                return text.strip()
        
        # Word DOCX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = docx.Document(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
            return text.strip()
        
        # PowerPoint PPTX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            pres = pptx.Presentation(file_path)
            text = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text).strip()
        
        # Excel XLSX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))
            return '\n'.join(text).strip()
        
        # Plain text
        elif mime_type in ['text/plain', 'text/csv', 'application/json']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read().strip()
        
        # Fallback: try to read as text
        else:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read().strip()
            except:
                logger.warning(f"Unsupported MIME type {mime_type} for {file_path}")
                return None
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        logger.debug(traceback.format_exc())
        return None

def update_database(db_path, file_path, raw_content):
    """Update raw_content column for given file_path."""
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    try:
        cursor.execute("UPDATE documents SET raw_content = ? WHERE file_path = ?", (raw_content, file_path))
        conn.commit()
        if cursor.rowcount == 0:
            logger.warning(f"No matching record for {file_path}")
        else:
            logger.info(f"Updated {file_path}")
    except Exception as e:
        logger.error(f"Database update failed for {file_path}: {e}")
        conn.rollback()
    finally:
        conn.close()

def main():
    # Check if raw_content column exists
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("PRAGMA table_info(documents)")
    columns = [col[1] for col in cursor.fetchall()]
    if 'raw_content' not in columns:
        logger.error("raw_content column missing in documents table")
        return
    conn.close()
    
    # Count total and processed
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM documents WHERE raw_content IS NOT NULL")
    already_processed = cursor.fetchone()[0]
    cursor.execute("SELECT COUNT(*) FROM documents")
    total = cursor.fetchone()[0]
    conn.close()
    logger.info(f"Already processed: {already_processed}/{total}")
    
    # Process each line in catalogue
    processed = 0
    errors = 0
    with open(CATALOGUE_PATH, 'r') as f:
        for line_num, line in enumerate(f, 1):
            line = line.strip()
            if not line:
                continue
            parts = line.split('\t')
            if len(parts) < 3:
                logger.warning(f"Invalid line {line_num}: {line}")
                continue
            file_path, size_str, mime_type = parts[0], parts[1], parts[2]
            
            # Skip if already processed
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            cursor.execute("SELECT raw_content FROM documents WHERE file_path = ?", (file_path,))
            row = cursor.fetchone()
            conn.close()
            if row and row[0] is not None:
                continue  # already has content
            
            # Extract text
            raw_content = extract_text(file_path, mime_type)
            if raw_content is None:
                logger.warning(f"Failed to extract text from {file_path}")
                errors += 1
                continue
            
            # Update database
            update_database(DB_PATH, file_path, raw_content)
            processed += 1
            
            if processed % 100 == 0:
                logger.info(f"Progress: {processed} files processed")
    
    logger.info(f"Done. Processed {processed} new files, {errors} errors.")

if __name__ == '__main__':
    main()